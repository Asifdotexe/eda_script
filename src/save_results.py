from pptx import Presentation
from pptx.util import Inches
import os

def create_ppt(output_directory: str, presentation_path: str) -> None:
    """This function creates a PowerPoint presentation from the given output directory.
    It reads the statistical analysis and chart files from the specified directories,
    adds them as slides in the presentation, and saves the final presentation at the specified path.

    :param output_directory: The directory containing the statistical analysis and chart files.
    :param presentation_path: The path where the final presentation will be saved.

    :returns: The function does not return any value, but it saves the presentation at the specified path.
    :rtype: None
    """
    prs = Presentation()

    # Get statistics and charts file lists
    stats_files = os.listdir(f"{output_directory}/stats/")
    charts_files = os.listdir(f"{output_directory}/charts/")

    # Add a slide for each statistical analysis file
    for stats_file in stats_files:
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide = prs.slides.add_slide(slide_layout)
        with open(f"{output_directory}/stats/{stats_file}") as file:
            stats_content = file.read()
        slide.shapes.title.text = f"Statistical Analysis - {stats_file}"
        slide.shapes.add_textbox(
            Inches(1),
            Inches(1.5),
            Inches(8),
            Inches(5)
        ).text = stats_content

    # Add a slide for each chart file
    for chart_file in charts_files:
        slide_layout = prs.slide_layouts[5]  # Title only layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Chart - {chart_file}"
        slide.shapes.add_picture(
            f"{output_directory}/charts/{chart_file}",
            Inches(1),
            Inches(1.5),
            Inches(8),
            Inches(5),
        )

    # Save the PowerPoint presentation
    prs.save(presentation_path)

def save_results(output_directory: str) -> str:
    """
    This function creates a PowerPoint presentation from the given output directory.
    It reads the statistical analysis and chart files from the specified directories,
    adds them as slides in the presentation, and saves the final presentation at the specified path.

    :param output_directory: The directory containing the statistical analysis and chart files.

    :returns: presentation_path: The path of the saved presentation file.
    """
    presentation_path = f"{output_directory}/eda_presentation.pptx"
    create_ppt(output_directory, presentation_path)
    return presentation_path